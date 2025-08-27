#!/usr/bin/env python3
# 0_batch_extract_txt_from_pdf_office_etc.py
# Extract markdown text from PDF and Office documents
#
# Supports: PDF, DOCX, PPTX, TXT, MD files
# Output: Converts files to markdown format with .pdf.md.txt, .docx.md.txt extensions
#
# Usage:
#   python 0_batch_extract_txt_from_pdf_office_etc.py                    # Use default directories
#   python 0_batch_extract_txt_from_pdf_office_etc.py <source> <target>  # Custom directories
#   python 0_batch_extract_txt_from_pdf_office_etc.py --help             # Show help

import sys
import pathlib
import subprocess
import logging
from datetime import datetime
from tqdm import tqdm

# Set up comprehensive logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('pdf_extraction_debug.log'),
        logging.StreamHandler()
    ]
)

# Supported file extensions
SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt", ".md"}

def print_help():
    """Print usage help."""
    print("""
PDF and Office Document Text Extraction

Usage:
  python 0_batch_extract_txt_from_pdf_office_etc.py [<source> <target>]

Arguments:
  <source>    Source directory with PDF/Office documents
  <target>    Target directory for extracted text files

Default Directories:
  Source: data/0_pdf_office_etc_source/
  Target: data/3_txt_transcribed/

Supported Formats:
  - PDF files (.pdf)
  - Word documents (.docx)
  - PowerPoint presentations (.pptx)
  - Text files (.txt)
  - Markdown files (.md)

Output Format:
  input.pdf        -> input.pdf.md.txt
  document.docx    -> document.docx.md.txt
  slides.pptx      -> slides.pptx.md.txt
  readme.txt       -> readme.txt (copied directly, collision backup as .old)
  notes.md         -> notes.md.txt

Requirements:
  - Python packages: pymupdf, python-docx, python-pptx, markdown
  - Install with: pip install pymupdf python-docx python-pptx markdown

Examples:
  python 0_batch_extract_txt_from_pdf_office_etc.py
  python 0_batch_extract_txt_from_pdf_office_etc.py documents/ extracted/
    """)

def check_dependencies():
    """Check if required Python packages are available."""
    required_packages = [
        ("fitz", "PyMuPDF for PDF processing"),
        ("docx", "python-docx for Word documents"),
        ("pptx", "python-pptx for PowerPoint files"),
        ("markdown", "markdown for Markdown processing")
    ]
    
    missing_packages = []
    for package_name, description in required_packages:
        try:
            __import__(package_name)
        except ImportError:
            missing_packages.append(f"{package_name} ({description})")
    
    if missing_packages:
        print("❌ Missing required Python packages:")
        for package in missing_packages:
            print(f"   - {package}")
        print()
        print("Install with: pip install pymupdf python-docx python-pptx markdown")
        sys.exit(1)
    
    print("✅ Dependencies check passed")

def extract_pdf_text(file_path: pathlib.Path) -> str:
    """Extract text from PDF using PyMuPDF."""
    import fitz  # PyMuPDF
    
    try:
        doc = fitz.open(file_path)
        text_content = []
        
        for page_num in range(doc.page_count):
            page = doc[page_num]
            text = page.get_text()
            if text.strip():
                text_content.append(f"# Page {page_num + 1}\n\n{text}\n")
        
        doc.close()
        return "\n".join(text_content)
    except Exception as e:
        return f"Error extracting PDF: {e}"

def extract_docx_text(file_path: pathlib.Path) -> str:
    """Extract text from DOCX using python-docx."""
    from docx import Document
    
    try:
        doc = Document(file_path)
        text_content = []
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                # Simple markdown formatting
                text = paragraph.text
                if paragraph.style.name.startswith('Heading'):
                    level = paragraph.style.name.split()[-1]
                    if level.isdigit():
                        text = f"{'#' * int(level)} {text}"
                text_content.append(text)
        
        return "\n\n".join(text_content)
    except Exception as e:
        return f"Error extracting DOCX: {e}"

def extract_pptx_text(file_path: pathlib.Path) -> str:
    """Extract text from PPTX using python-pptx."""
    from pptx import Presentation
    
    try:
        prs = Presentation(file_path)
        text_content = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = f"# Slide {slide_num}\n\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text += f"{shape.text}\n\n"
            
            if slide_text != f"# Slide {slide_num}\n\n":
                text_content.append(slide_text)
        
        return "\n".join(text_content)
    except Exception as e:
        return f"Error extracting PPTX: {e}"

def extract_txt_text(file_path: pathlib.Path) -> str:
    """Read plain text file."""
    try:
        return file_path.read_text(encoding='utf-8')
    except UnicodeDecodeError:
        try:
            return file_path.read_text(encoding='latin-1')
        except Exception as e:
            return f"Error reading text file: {e}"

def extract_md_text(file_path: pathlib.Path) -> str:
    """Read markdown file and optionally convert to formatted text."""
    try:
        content = file_path.read_text(encoding='utf-8')
        # For markdown files, we'll just return the raw markdown
        # Could optionally convert to HTML then to text using markdown package
        return content
    except Exception as e:
        return f"Error reading markdown file: {e}"

def extract_text_from_file(file_path: pathlib.Path) -> str:
    """Extract text from supported file formats."""
    extension = file_path.suffix.lower()
    
    if extension == ".pdf":
        return extract_pdf_text(file_path)
    elif extension == ".docx":
        return extract_docx_text(file_path)
    elif extension == ".pptx":
        return extract_pptx_text(file_path)
    elif extension == ".txt":
        return extract_txt_text(file_path)
    elif extension == ".md":
        return extract_md_text(file_path)
    else:
        return f"Unsupported file format: {extension}"

def handle_txt_collision(target_path: pathlib.Path):
    """Handle collision for TXT files by creating .old backups."""
    if not target_path.exists():
        return target_path
    
    # Create backup chain: .old, .old.old, .old.old.old, etc.
    current_path = target_path
    while current_path.exists():
        backup_path = pathlib.Path(str(current_path) + ".old")
        if not backup_path.exists():
            current_path.rename(backup_path)
            break
        current_path = backup_path
    
    return target_path

def process_files(source_dir: pathlib.Path, target_dir: pathlib.Path):
    """Process all supported files in source directory."""
    logging.info(f"Starting PDF/Office extraction process")
    logging.info(f"Source directory: {source_dir}")
    logging.info(f"Target directory: {target_dir}")
    
    # Check if target directory exists and what files are in it
    if target_dir.exists():
        existing_files = list(target_dir.glob("*"))
        logging.info(f"Target directory exists with {len(existing_files)} existing files:")
        for file in existing_files:
            logging.info(f"  Existing: {file.name} ({file.stat().st_size} bytes)")
    else:
        logging.info(f"Target directory does not exist, will create it")
    
    # Ensure target directory exists
    target_dir.mkdir(parents=True, exist_ok=True)
    logging.info(f"Target directory created/confirmed: {target_dir}")
    
    # Find all supported files
    supported_files = []
    for ext in SUPPORTED_EXTENSIONS:
        supported_files.extend(source_dir.glob(f"*{ext}"))
    
    if not supported_files:
        print(f"No supported files found in {source_dir}")
        print(f"Supported formats: {', '.join(sorted(SUPPORTED_EXTENSIONS))}")
        return
    
    print(f"Found {len(supported_files)} supported files")
    
    # Process each file
    for file_path in tqdm(supported_files, desc="Processing files", unit="file"):
        try:
            logging.info(f"Processing file: {file_path.name} ({file_path.suffix.lower()})")
            extension = file_path.suffix.lower()
            
            if extension == ".txt":
                # TXT files: copy directly with collision handling
                target_path = target_dir / file_path.name
                logging.info(f"TXT file - target path: {target_path}")
                target_path = handle_txt_collision(target_path)
                logging.info(f"TXT file - final target path after collision check: {target_path}")
                
                # Read and copy content
                content = extract_txt_text(file_path)
                target_path.write_text(content, encoding='utf-8')
                logging.info(f"Successfully created TXT file: {target_path} ({len(content)} chars)")
                
            elif extension == ".md":
                # MD files: copy content to .md.txt
                output_filename = f"{file_path.stem}.md.txt"
                output_path = target_dir / output_filename
                logging.info(f"MD file - output path: {output_path}")
                
                content = extract_md_text(file_path)
                output_path.write_text(content, encoding='utf-8')
                logging.info(f"Successfully created MD file: {output_path} ({len(content)} chars)")
                
            else:
                # PDF, DOCX, PPTX: extract text to .ext.md.txt
                logging.info(f"Extracting text from {extension.upper()} file: {file_path.name}")
                extracted_text = extract_text_from_file(file_path)
                
                output_filename = f"{file_path.name}.md.txt"
                output_path = target_dir / output_filename
                logging.info(f"Output path for {extension.upper()}: {output_path}")
                
                output_path.write_text(extracted_text, encoding='utf-8')
                
                # Verify file was created and log details
                if output_path.exists():
                    file_size = output_path.stat().st_size
                    logging.info(f"Successfully created {extension.upper()} extraction file: {output_path} ({file_size} bytes, {len(extracted_text)} chars)")
                else:
                    logging.error(f"Failed to create output file: {output_path}")
            
        except Exception as e:
            logging.error(f"Error processing {file_path.name}: {e}")
            # Create error file
            error_filename = f"{file_path.name}.error.txt"
            error_path = target_dir / error_filename
            error_path.write_text(f"Error processing {file_path.name}: {e}\n", encoding='utf-8')
            logging.info(f"Created error file: {error_path}")

def main():
    start_time = datetime.now()
    logging.info(f"=== PDF/Office Text Extraction Started at {start_time} ===")
    # Parse arguments
    args = sys.argv[1:]
    
    # Handle --help
    if "--help" in args or "-h" in args:
        print_help()
        sys.exit(0)
    
    # Parse directory arguments
    if len(args) == 0:
        # Use default directories
        source_dir = pathlib.Path("data/0_pdf_office_etc_source").expanduser()
        target_dir = pathlib.Path("data/3_txt_transcribed").expanduser()
        print(f"Using default directories: {source_dir} -> {target_dir}")
    elif len(args) == 2:
        source_dir = pathlib.Path(args[0]).expanduser()
        target_dir = pathlib.Path(args[1]).expanduser()
        print(f"Using custom directories: {source_dir} -> {target_dir}")
    else:
        print("❌ Error: Invalid number of arguments")
        print()
        print("Usage:")
        print("  python 0_batch_extract_txt_from_pdf_office_etc.py                    # Default directories")
        print("  python 0_batch_extract_txt_from_pdf_office_etc.py <source> <target>  # Custom directories")
        print()
        print("Use --help for detailed information")
        sys.exit(1)
    
    # Validate source directory
    if not source_dir.exists():
        print(f"❌ Error: Source directory not found: {source_dir}")
        print("Create the directory and add PDF/Office documents to process")
        sys.exit(1)
    
    # Check dependencies
    check_dependencies()
    
    print(f"🚀 Starting PDF and Office document text extraction")
    print(f"   Source: {source_dir}")
    print(f"   Target: {target_dir}")
    
    # Log final state before processing
    logging.info(f"About to process files from {source_dir} to {target_dir}")
    
    # Process files
    process_files(source_dir, target_dir)
    
    # Log final state after processing
    if target_dir.exists():
        final_files = list(target_dir.glob("*"))
        logging.info(f"Processing complete. Final target directory contains {len(final_files)} files:")
        for file in final_files:
            size = file.stat().st_size
            logging.info(f"  Final: {file.name} ({size} bytes)")
    else:
        logging.error(f"Target directory {target_dir} does not exist after processing!")
    
    print(f"✅ Text extraction complete!")
    print(f"   Output files saved to: {target_dir}")
    print(f"   Files ready for further processing:")
    print(f"   - PDF/DOCX/PPTX: converted to .pdf.md.txt, .docx.md.txt, .pptx.md.txt")
    print(f"   - MD files: copied to .md.txt")
    print(f"   - TXT files: copied directly (with collision backup)")
    print(f"   - Full processing log saved to: pdf_extraction_debug.log")
    logging.info(f"PDF/Office extraction process completed successfully at {datetime.now()}")

if __name__ == "__main__":
    main()